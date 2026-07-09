"""Multi-format text extraction, dispatched by file extension.

Lean by design: `.txt`/`.md`/`.html`/`.htm`/`.vtt`/`.srt` are native (stdlib
only, no optional deps). Richer document formats (`.pdf`/`.docx`/`.pptx`/
`.xlsx`) lazily import their library *inside* the extractor function so the
core `htc` install stays dependency-light — install the `ingest` extra
(`pip install htc[ingest]`) to enable them.
"""

from __future__ import annotations

import re
from html.parser import HTMLParser
from pathlib import Path

_INGEST_EXTRA_HINT = "requires the 'ingest' extra: pip install htc[ingest]"


class UnsupportedFormatError(ValueError):
    """Raised when a file extension has no registered extractor."""


class MissingDependencyError(ImportError):
    """Raised when an optional extraction dependency isn't installed."""


class _HTMLTextExtractor(HTMLParser):
    """Strips tags, keeping only text nodes (script/style content excluded)."""

    def __init__(self) -> None:
        super().__init__()
        self._chunks: list[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in ("script", "style"):
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style") and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0:
            self._chunks.append(data)

    def text(self) -> str:
        return "".join(self._chunks)


def _extract_html(path: Path) -> str:
    parser = _HTMLTextExtractor()
    parser.feed(path.read_text(errors="replace"))
    return parser.text()


_TIMESTAMP_LINE = re.compile(r"^\d{2}:\d{2}:\d{2}[.,]\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}[.,]\d{3}.*$")
_CUE_INDEX_LINE = re.compile(r"^\d+$")


def _extract_transcript(path: Path) -> str:
    """Strip WebVTT/SRT timestamp + cue-index lines, keep spoken text."""
    lines = path.read_text(errors="replace").splitlines()
    kept = []
    for line in lines:
        stripped = line.strip()
        if not stripped or stripped.upper() == "WEBVTT":
            continue
        if _TIMESTAMP_LINE.match(stripped) or _CUE_INDEX_LINE.match(stripped):
            continue
        kept.append(stripped)
    return "\n".join(kept)


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as err:
        raise MissingDependencyError(f"extracting .pdf {_INGEST_EXTRA_HINT}") from err
    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(path: Path) -> str:
    try:
        import docx
    except ImportError as err:
        raise MissingDependencyError(f"extracting .docx {_INGEST_EXTRA_HINT}") from err
    document = docx.Document(str(path))
    return "\n\n".join(p.text for p in document.paragraphs)


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as err:
        raise MissingDependencyError(f"extracting .pptx {_INGEST_EXTRA_HINT}") from err
    presentation = Presentation(str(path))
    slides_text = []
    for slide in presentation.slides:
        frames = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        slides_text.append("\n".join(frames))
    return "\n\n".join(slides_text)


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as err:
        raise MissingDependencyError(f"extracting .xlsx {_INGEST_EXTRA_HINT}") from err
    workbook = load_workbook(str(path), read_only=True, data_only=True)
    rows_text = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows_text.append("\t".join(cells))
    return "\n".join(rows_text)


def _extract_plain(path: Path) -> str:
    return path.read_text(errors="replace")


# Common source-code suffixes — extracted as plaintext, no parsing. This is
# what makes `ingest_sources` cover code as well as docs, so memory and the
# knowledge graph both ground on the actual implementation, not just docs.
CODE_SUFFIXES = frozenset(
    {
        ".py",
        ".js",
        ".ts",
        ".tsx",
        ".jsx",
        ".go",
        ".rs",
        ".java",
        ".rb",
        ".php",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".cs",
        ".swift",
        ".kt",
        ".scala",
        ".sh",
        ".css",
        ".scss",
        ".sql",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
    }
)

_EXTRACTORS = {
    ".txt": _extract_plain,
    ".md": _extract_plain,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".vtt": _extract_transcript,
    ".srt": _extract_transcript,
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    **{suffix: _extract_plain for suffix in CODE_SUFFIXES},
}

# Extensions with a registered extractor — used by callers that want to
# pre-filter files without invoking extraction.
KNOWN_SUFFIXES = frozenset(_EXTRACTORS)

# Mirrors `goldens.generator.SECRET_SUFFIXES`/`SECRET_NAMES`/`SECRET_STEMS`:
# duplicated (not imported) to avoid a circular import (goldens.generator
# already imports `world_model.ingest`). Keep these two lists in sync.
_SECRET_SUFFIXES = {".pem", ".key", ".p12", ".pfx", ".keystore", ".jks"}
_SECRET_NAMES = {"credentials", "secrets", ".npmrc", ".pypirc", ".netrc", ".htpasswd"}
_SECRET_STEMS = ("id_rsa", "id_ed25519", "id_dsa", "id_ecdsa")


def is_secret_file(path: Path) -> bool:
    """True if `path` looks like it carries secrets (.env*, keys, credentials)
    — such files must never enter the ingest corpus."""
    name = path.name.lower()
    if name.startswith(".env"):
        return True
    if path.suffix.lower() in _SECRET_SUFFIXES:
        return True
    if name in _SECRET_NAMES:
        return True
    return any(name.startswith(stem) for stem in _SECRET_STEMS)


def extract_text(path: Path) -> str:
    """Extract plain text from `path`, dispatching by extension.

    Raises `UnsupportedFormatError` for unregistered extensions and
    `MissingDependencyError` when the format needs an optional dependency
    that isn't installed (`pip install htc[ingest]`).
    """
    extractor = _EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        raise UnsupportedFormatError(f"no extractor for extension: {path.suffix or '(none)'}")
    return extractor(path)
